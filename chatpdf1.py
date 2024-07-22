import streamlit as st
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from dotenv import load_dotenv
import os

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Verify the API key is loaded
print("Google API Key:", os.getenv("GOOGLE_API_KEY"))

def extract_text_from_pdf(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(docx_docs):
    text = ""
    for doc in docx_docs:
        document = Document(doc)
        for para in document.paragraphs:
            text += para.text + "\n"
    return text

def extract_text_from_pptx(pptx_docs):
    text = ""
    for ppt in pptx_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
    chunks = text_splitter.split_text(text)
    return chunks

def get_summary(text):
    summary_prompt = """
    Summarize the following text into a concise summary:\n\n
    {context}
    """
    
    prompt = PromptTemplate(template=summary_prompt, input_variables=["context"])
    model = ChatGoogleGenerativeAI(model="models/gemini-pro", temperature=0.3)
    chain = LLMChain(prompt=prompt, llm=model)
    
    try:
        summary = chain.run({"context": text})
        return summary
    except Exception as e:
        st.error(f"Error generating summary: {e}")
        return "Error generating summary."

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", api_key=os.getenv("GOOGLE_API_KEY"))
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context, make sure to provide all the details, if the answer is not in
    provided context just say, "answer is not available in the context", don't provide the wrong answer\n\n
    Context:\n {context}\n
    Question: \n{question}\n

    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3, api_key=os.getenv("GOOGLE_API_KEY"))
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001", api_key=os.getenv("GOOGLE_API_KEY"))
        new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
        chain = get_conversational_chain()
        response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
        st.session_state.chat_history.append({"question": user_question, "response": response["output_text"]})
        st.session_state.chat_history_str += f"User: {user_question}\nBot: {response['output_text']}\n\n"
    except Exception as e:
        st.error(f"Error: {e}")

def display_chat_history():
    if "chat_history_str" in st.session_state:
        st.text_area("Chat History", value=st.session_state.chat_history_str, height=400)

def main():
    st.set_page_config(page_title="Chat with InfoFlux", page_icon=":speech_balloon:")
    st.header("Chat with InfoFlux💁")

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
        st.session_state.chat_history_str = ""

    with st.sidebar:
        st.title("Upload Documents:")
        pdf_docs = st.file_uploader("Upload PDF Files", type="pdf", accept_multiple_files=True)
        docx_docs = st.file_uploader("Upload Word Documents", type="docx", accept_multiple_files=True)
        pptx_docs = st.file_uploader("Upload PowerPoint Presentations", type="pptx", accept_multiple_files=True)

        if st.button("Submit & Process"):
            with st.spinner("Processing..."):
                raw_text = ""
                if pdf_docs:
                    raw_text += extract_text_from_pdf(pdf_docs)
                if docx_docs:
                    raw_text += extract_text_from_docx(docx_docs)
                if pptx_docs:
                    raw_text += extract_text_from_pptx(pptx_docs)

                if raw_text:
                    text_chunks = get_text_chunks(raw_text)
                    summarized_chunks = [get_summary(chunk) for chunk in text_chunks]
                    get_vector_store(summarized_chunks)
                    st.success("Processing Complete")
                else:
                    st.warning("No documents uploaded")

    # User input section
    user_question = st.text_input("Ask a Question from the Documents")

    if user_question:
        user_input(user_question)

    # Display chat history
    display_chat_history()

if __name__ == "__main__":
    main()

